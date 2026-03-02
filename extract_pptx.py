"""
Script to extract text and images from row.pptx
"""
import os
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Paths
PPTX_FILE = r"C:\Users\Ajitha.Rajkumar\chatbot\row\row.pptx"
OUTPUT_DIR = r"C:\Users\Ajitha.Rajkumar\chatbot\row\extracted"
IMAGES_DIR = os.path.join(OUTPUT_DIR, "images")
TEXT_FILE = os.path.join(OUTPUT_DIR, "extracted_text.txt")

# Create output directories
os.makedirs(IMAGES_DIR, exist_ok=True)

# Open presentation
prs = Presentation(PPTX_FILE)

image_count = 0
text_lines = []

for slide_num, slide in enumerate(prs.slides, start=1):
    slide_header = f"\n{'='*60}\nSlide {slide_num}\n{'='*60}"
    text_lines.append(slide_header)
    print(slide_header)

    for shape in slide.shapes:
        # --- Extract text ---
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                para_text = paragraph.text.strip()
                if para_text:
                    text_lines.append(para_text)
                    print(f"  [Text] {para_text}")

        # --- Extract table text ---
        if shape.has_table:
            table = shape.table
            text_lines.append("\n  [Table]")
            print("  [Table]")
            for row_idx, row in enumerate(table.rows):
                row_data = []
                for cell in row.cells:
                    row_data.append(cell.text.strip())
                row_str = " | ".join(row_data)
                text_lines.append(f"    Row {row_idx + 1}: {row_str}")
                print(f"    Row {row_idx + 1}: {row_str}")

        # --- Extract images ---
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_count += 1
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            image_filename = f"slide{slide_num}_image{image_count}.{ext}"
            image_path = os.path.join(IMAGES_DIR, image_filename)
            with open(image_path, "wb") as f:
                f.write(image.blob)
            text_lines.append(f"  [Image extracted: {image_filename}]")
            print(f"  [Image] Saved: {image_filename}")

        # --- Extract images from group shapes ---
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for grp_shape in shape.shapes:
                if grp_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image_count += 1
                    image = grp_shape.image
                    ext = image.content_type.split("/")[-1]
                    if ext == "jpeg":
                        ext = "jpg"
                    image_filename = f"slide{slide_num}_image{image_count}.{ext}"
                    image_path = os.path.join(IMAGES_DIR, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(image.blob)
                    text_lines.append(f"  [Image extracted: {image_filename}]")
                    print(f"  [Image] Saved: {image_filename}")

# Write all extracted text to a file
with open(TEXT_FILE, "w", encoding="utf-8") as f:
    f.write("\n".join(text_lines))

print(f"\n{'='*60}")
print(f"Extraction complete!")
print(f"  Total slides: {len(prs.slides)}")
print(f"  Total images extracted: {image_count}")
print(f"  Text saved to: {TEXT_FILE}")
print(f"  Images saved to: {IMAGES_DIR}")
print(f"{'='*60}")
